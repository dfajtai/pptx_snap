import os.path
from abc import abstractmethod
from typing import Optional, Any
import numpy as np

from pptx.util import Length, Cm, Mm

from .pptx_reader import PPTXReader
from .grid import Grid
from .snappable_object import SnappableObject
from .slide import Slide
from .templates import ObjectTemplates
from .utils import AnchorPoint

class SnapCandidate:
    """
    Class to store relevant information on 'SnapCandidates'
    A SnapCandidate is a virtual, grid-sanpped position of an objects anchor point
    """

    def __init__(self, snappable_object: SnappableObject, anchor_point: AnchorPoint, snap_type:str, grid_type:str, snap_x_position: Optional[int] = None, snap_y_position: Optional[int] = None):
        self.snappable_object = snappable_object
        self.anchor_point = anchor_point      # Name of the anchor point (e.g., "top-left")
        self.snap_type = snap_type
        self.grid_type = grid_type
        
        self.anchor_position = snappable_object.get_anchor_point(anchor_point)
        
        x = snap_x_position if not isinstance(snap_x_position,type(None)) else self.anchor_position[0]
        y = snap_y_position if not isinstance(snap_y_position,type(None)) else self.anchor_position[1]
        
        self.snap_position = (x,y)    # New snapped position (x, y)

    def __str__(self) -> str:
        return (f"{self.anchor_point.value} anchor at {self.anchor_position} "
                f"snapped to {self.snap_position}")
        
    @property
    def displacement_vector(self)-> np.ndarray:
        return np.array(self.snap_position) - np.array(self.anchor_position)

    @property
    def relative_displacement_vector(self) -> np.ndarray:
        return  np.abs(self.displacement_vector)/np.array(self.snappable_object.sizes)

    @property
    def displacement(self)-> float:
        return float(np.linalg.norm(self.displacement_vector))



class Snapping:
    """
    Base class to calculate SnapCandidates
    """

    def __init__(self, grid: Grid) -> None:
        self.grid = grid
        
    @abstractmethod
    def snap(self,obj:SnappableObject, grid_type:str) -> None:
        pass


class XSnapping(Snapping):
    """
    Class to calculate SnapCandidates on X axis
    """

    def __init__(self, grid: Grid):
        super().__init__(grid)
        self.snap_type = "x"

    def snap(self, obj: SnappableObject, grid_type:str):
        """Apply x-axis snapping to the object."""
        for anchor_point in obj.active_anchor_points:
            if len(self.grid.x_grid_lines)==0:
                break
            closest_x = min(self.grid.x_grid_lines, key=lambda gx: abs(gx - obj.get_anchor_point(anchor_point)[0]))
            obj.snapping_candidates.append(SnapCandidate(obj, anchor_point = anchor_point, snap_x_position=closest_x,
                                                         snap_type = self.snap_type,
                                                         grid_type=grid_type))


class YSnapping(Snapping):
    """
    Class to calculate SnapCandidates on Y axis
    """

    def __init__(self, grid: Grid) -> None:
        super().__init__(grid)
        self.snap_type = "y"

    def snap(self, obj: SnappableObject, grid_type:str) -> None:
        """Apply y-axis snapping to the object."""
        for anchor_point in obj.active_anchor_points:
            if len(self.grid.y_grid_lines)==0:
                break
            closest_y = min(self.grid.y_grid_lines, key=lambda gy: abs(gy - obj.get_anchor_point(anchor_point)[1]))
            obj.snapping_candidates.append(SnapCandidate(obj, anchor_point=anchor_point, snap_y_position=closest_y,
                                                         snap_type = self.snap_type,
                                                         grid_type=grid_type))


class JointSnapping(Snapping):
    """
    Class to calculate SnapCandidates on X and Y axes
    """

    def __init__(self, grid: Grid) -> None:
        super().__init__(grid)
        self.snap_type = "joint"

    def snap(self, obj: SnappableObject, grid_type:str) -> None:
        """Apply simultaneous x and y snapping to the object."""
        for anchor_point in obj.active_anchor_points:
            if len(self.grid.x_grid_lines) != 0:
                closest_x = min(self.grid.x_grid_lines, key=lambda gx: abs(gx - obj.get_anchor_point(anchor_point)[0]))
            else:
                closest_x = obj.get_anchor_point(anchor_point)[0]

            if len(self.grid.y_grid_lines) != 0:
                closest_y = min(self.grid.y_grid_lines, key=lambda gy: abs(gy - obj.get_anchor_point(anchor_point)[1]))
            else:
                closest_y = obj.get_anchor_point(anchor_point)[1]

            obj.snapping_candidates.append(SnapCandidate(obj, anchor_point=anchor_point, snap_x_position=closest_x, snap_y_position=closest_y,
                                                         snap_type = self.snap_type,
                                                         grid_type=grid_type))
        
        
class SnappingSearch:
    """
    Class implement methods to calculate various SnappingCandidates for SnappableObjects
    """

    def __init__(self, allow_x_snap = True, allow_y_snap = True):  

        self.allow_x_snap = allow_x_snap
        self.allow_y_snap = allow_y_snap
        
        self.x_grid = None
        self.y_grid = None
        self.joint_grid = None
        
        self.snapping_strategies = {}


        
    def _update_strategies(self)-> None:
        self.snapping_strategies = {}
        if isinstance(self.x_grid,Grid) and self.allow_x_snap:
            self.snapping_strategies["x"] = XSnapping(self.x_grid)
            
        if isinstance(self.y_grid,Grid) and self.allow_y_snap:
            self.snapping_strategies["y"] = XSnapping(self.y_grid)
            
        if isinstance(self.x_grid,Grid) and isinstance(self.y_grid,Grid) and self.allow_x_snap and self.allow_y_snap:
            x_grid = self.x_grid.get_x_grid()
            y_grid = self.y_grid.get_y_grid()
            
            self.joint_grid = Grid.merge_grids(x_grid, y_grid)
            
            self.snapping_strategies["joint"] = JointSnapping(self.joint_grid)
    
    
    def set_joint_grid(self,grid:Grid) -> None:
        assert isinstance(grid,Grid)
        self.x_grid = grid.get_x_grid()
        self.y_grid = grid.get_y_grid()
        self._update_strategies()
            
    def set_x_grid(self,grid:Grid) -> None:
        assert isinstance(grid,Grid)
        self.x_grid = grid
        self._update_strategies()
    
        
    def set_y_grid(self, grid:Grid)-> None:
        assert isinstance(grid,Grid)
        self.y_grid = grid
        self._update_strategies()
    
        
    def extend_x_grid(self, grid:Grid)-> None:
        assert isinstance(grid,Grid)
        self.x_grid.extend(grid)
        self._update_strategies()
    
    
    def extend_y_grid(self, grid:Grid)-> None:
        assert isinstance(grid,Grid)
        self.y_grid.extend(grid)
        self._update_strategies()
    
    
    def calculate_candidates_for_all_obj(self, slide:Slide, strategy_type: str, flush = False, grid_type:str = "unknown"):
        """Apply the given snapping strategy (x, y, or joint) for all SnappableObject on a given Slide"""
        assert isinstance(slide,Slide)
        
        for so in slide.snappable_objects:
            self.calculate_candidates(so,strategy_type=strategy_type,flush=flush, grid_type= grid_type)
            
    
    def calculate_candidates(self, obj: SnappableObject, strategy_type: str, flush = False, grid_type:str = "unknown"):
        """Apply the given snapping strategy (x, y, or joint) for a given SnappableObject"""
        assert isinstance(obj,SnappableObject)
        
        strategy = self.snapping_strategies.get(strategy_type)
        if strategy:
            if isinstance(strategy,Snapping):
                if flush:
                    obj.snapping_candidates.clear()
                    
                strategy.snap(obj,grid_type=grid_type)
    

class SnappingManager:
    """
    Class to search and manifest the best snap candidate for a SnappableObject
    """

    def __init__(self,
                 reader: PPTXReader,
                 x_limit:Optional[Length] = None,
                 y_limit:Optional[Length] = None,
                 x_relative_limit: Optional[float] = None,
                 y_relative_limit: Optional[float] = None)->None:

        self.reader = reader
        self.fix_limit = np.array([l.emu if isinstance(l,Length) else l for l in [x_limit,y_limit] ])
        self.rel_limit = np.array([x_relative_limit,y_relative_limit])


    def _validate_displacement(self, displacement_vector:np.ndarray) -> bool:
        if np.all(self.fix_limit == None): return True

        invalid_indices = np.isnan(self.fix_limit)
        return np.all(np.less_equal(displacement_vector[~invalid_indices],self.fix_limit[~invalid_indices]))

    def _validate_relative_displacement(self, rel_displacement_vector:np.ndarray) -> bool:
        if np.all(self.rel_limit == None): return True

        invalid_indices = np.isnan(self.rel_limit)
        return np.all(np.less_equal(rel_displacement_vector[~invalid_indices],self.rel_limit[~invalid_indices]))

    def apply_snaps(self):
        for slide_index, slide in enumerate(self.reader.slides):
            for so in slide.snappable_objects:
                valid_candidates = []
                for sc in so.snapping_candidates:
                    if not isinstance(sc,SnapCandidate):
                        continue

                    if not self._validate_displacement(sc.displacement_vector):
                        continue

                    if not self._validate_relative_displacement(sc.relative_displacement_vector):
                        continue

                    valid_candidates.append(sc)

                valid_candidates.sort(key=lambda x: x.displacement)

                if len(valid_candidates)==0:
                    continue

                best_candidate = valid_candidates[0]

                print(best_candidate)

                self.reader.presentation.slides[so.slide_index].shapes[so.shape_index].left +=best_candidate.displacement_vector[0]
                self.reader.presentation.slides[so.slide_index].shapes[so.shape_index].top +=best_candidate.displacement_vector[1]



    def save_at(self, out_path):
        if not os.path.exists(os.path.basename(out_path)):
            os.mkdir(os.path.basename(out_path))

        self.reader.presentation.save(out_path)