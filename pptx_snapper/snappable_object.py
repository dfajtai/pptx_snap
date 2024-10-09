from copy import deepcopy
from typing import Any
import numpy as np

import weakref

from collections import OrderedDict


from pptx.util import Length
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape

from .utils import AnchorPoint, classproperty


class SnappableObject:

    _default_active_anchor_points = [AnchorPoint.TOP_LEFT,
                                    AnchorPoint.TOP_RIGHT,
                                    AnchorPoint.BOTTOM_LEFT,
                                    AnchorPoint.BOTTOM_RIGHT,
                                    AnchorPoint.CENTER]

    catalog = weakref.WeakSet()
    template_catalog = weakref.WeakSet()

    def __init__(self, shape: BaseShape, slide_index: int, shape_index: int, is_template:bool = False):
        self.shape = shape
        self.slide_index = slide_index
        self.shape_index = shape_index

        self.is_template = is_template

        self.shape_id = shape.shape_id
        self.name = shape.name
        
        self.is_placeholder = shape.is_placeholder
        self.is_table = shape.has_table
        self.is_chart = shape.has_chart
        self.is_text = shape.has_text_frame
        self.is_picture = isinstance(shape,Picture)
        self.is_group = isinstance(shape,GroupShape)
        
        if self.is_text:
            self.text = shape.text

        self.left = shape.left
        self.top = shape.top
        self.width = shape.width
        self.height = shape.height

        self._left = shape.left
        self._top = shape.top
        self._width = shape.width
        self._height= shape.height

        self._template_snap_id = None  # Initially None, will be set later

        self._active_anchor_points = None
        self.snapping_candidates = []

        if is_template:
            self.__class__.template_catalog.add(self)
        else:
            self.__class__.catalog.add(self)

    @property
    def full_id(self) -> str:
        return f"{self.shape_id}#{self.shape_index}@{self.slide_index}"

    @property
    def sizes(self) -> np.ndarray:
        return np.array([Length(max(self.width,1)), Length(max(self.height,1))])

    @property
    def center(self)->tuple[Length,...]:
        return Length(self.left + self.width // 2), Length(self.top + self.height // 2)

    @property
    def right(self):
        return Length(self.left + self.width)

    @property
    def bottom(self):
        return Length(self.top + self.height)

    @property
    def corners(self)-> list[tuple[Length,...]]:
        """Calculate the corner points of the object."""
        return [
            (self.left, self.top),  # Top-left
            (self.right, self.top),  # Top-right
            (self.left, self.bottom),  # Bottom-left
            (self.right, self.bottom),  # Bottom-right
        ]

    @property
    def anchor_points(self)->dict[AnchorPoint,tuple[Length,...]]:
        corners = self.corners
        
        return OrderedDict({
            AnchorPoint.TOP_LEFT:corners[0],
            AnchorPoint.TOP_RIGHT: corners[1],
            AnchorPoint.BOTTOM_LEFT: corners[2],
            AnchorPoint.BOTTOM_RIGHT: corners[3],
            AnchorPoint.CENTER: self.center,
        })


    def get_anchor_point(self, anchor_point: AnchorPoint) -> tuple[Length, ...]:
        """Get the position of a named anchor point."""
        return self.anchor_points.get(anchor_point, (None, None))

    def get_anchor_point_by_name(self, anchor_point_name: str) -> tuple[Length, ...]:
        """Get the position of a named anchor point."""
        anchor_point = AnchorPoint.__getitem__(anchor_point_name)
        return self.get_anchor_point(anchor_point)

    @property
    def active_anchor_points(self)-> list[AnchorPoint]:
        if self._active_anchor_points:
            return self._active_anchor_points
        return SnappableObject._default_active_anchor_points

    @active_anchor_points.setter
    def active_anchor_points(self, new_anchor_points: list[AnchorPoint]) -> None:
        self._active_anchor_points = new_anchor_points[:]

    @classproperty
    def default_active_anchor_points(cls) -> list[AnchorPoint]:
        return cls._default_active_anchor_points

    @default_active_anchor_points.setter
    def default_active_anchor_points(cls,new_anchor_points: list[AnchorPoint]) -> None:
        cls._default_active_anchor_points = new_anchor_points[:]

    @classmethod
    def update_default_anchor_points(cls,new_anchor_points: list[AnchorPoint] | None = None, propagate = False) -> None:
        """
        Sets/updates the default anchor points for the class or its instances
        :param new_anchor_points:
        :param propagate: propagate the change to initialized instance variables ot not
        :return: nothing at all
        """

        if new_anchor_points:
            SnappableObject._default_active_anchor_points = new_anchor_points[:]

        if propagate:
            for obj in cls.__subclasses__():
                for instance in obj.__dict__.values():
                    if isinstance(instance, SnappableObject):
                        instance._active_anchor_points = None


    @property
    def area(self) -> int:
        return int(self.width * self.height)
    
    @property
    def orig_top(self) -> Length:
        return Length(self._top)

    @property
    def orig_left(self) -> Length:
        return Length(self._left)

    @property
    def orig_right(self) -> Length:
        return Length(self._left + self._width)

    @property
    def orig_bottom(self) -> Length:
        return Length(self._top + self._height)

    @property
    def orig_width(self) -> Length:
        return Length(self._width)

    @property
    def orig_height(self) -> Length:
        return Length(self._height)
    
    @property
    def shape_type(self) -> str:
        shape_type = "Shape"
        if self.is_picture: shape_type = "Picture"
        if self.is_table: shape_type = "Table"
        if self.is_chart: shape_type = "Chart"
        if self.is_text: shape_type = "Text"
        if self.is_group: shape_type = "Group"
        
        return shape_type


    @property
    def template_snap_id(self) -> str:
        return self._template_snap_id

    @template_snap_id.setter
    def template_snap_id(self, value:str):
        self._template_snap_id = value

    def intersection_area(self, other:'SnappableObject') -> float:
        x_left = max(self.left, other.left)
        y_top = max(self.top, other.top)
        x_right = min(self.right, other.right)
        y_bottom = min(self.bottom, other.bottom)

        overlap_width = max(0, x_right - x_left)
        overlap_height = max(0, y_bottom - y_top)

        return overlap_width * overlap_height

    
    def size_match_score(self, other: 'SnappableObject') -> float:
        width_diff = abs(self.width - other.width)
        height_diff = abs(self.height - other.height)
        score = 1.0 - (width_diff / max(self.width, other.width)) * (height_diff / max(self.height, other.height))
        return max(0, min(1, score))  # Clamp between 0 and 1
    
    
    def dice_coefficient(self, other: 'SnappableObject') -> float:
        intersection_area = self.intersection_area(other)
        if self.area + other.area == 0:
            return 0
        return (2 * intersection_area) / (self.area + other.area)
    
    def overlap_percentage(self, other: 'SnappableObject') -> float:
        intersection_area = self.intersection_area(other)
        min_area = min(self.area, other.area)
        if min_area == 0:
            return 0
        return (intersection_area / min_area) * 100
    
    
    def size_corrected_dice(self, other: 'SnappableObject') -> float:
        """Calculate the size-corrected Dice score between two SnappableObjects."""
        # Compute the area of each object
        area1 = self.area
        area2 = other.area

        # Compute the overlap area between the objects
        overlap_area = self.intersection_area(other)

        # Calculate the standard Dice coefficient
        dice_coefficient = (2 * overlap_area) / (area1 + area2)

        # Calculate the size correction factor
        size_correction = 1 - abs(area1 - area2) / max(area1, area2)

        # Compute the size-corrected Dice score
        size_corrected_dice = dice_coefficient * size_correction

        return size_corrected_dice

    
    def __str__(self) -> str:
        return f"On [Slide {self.slide_index}] Shape with type of '{self.shape_type}' called '{self.name}' @ [{self.left};{self.top}] with size of [{self.width} x {self.height}]"