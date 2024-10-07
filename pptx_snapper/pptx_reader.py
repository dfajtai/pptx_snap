from pptx import Presentation
from .slide import Slide

class PPTXReader:
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.presentation = Presentation(file_path)
        
        self.slide_width = self.presentation.slide_width
        self.slide_height = self.presentation.slide_height
    
        self.slides = self.read_slides()

    def read_slides(self):
        """Read all slides and store them as Slide objects."""
        slides = []
        for index, slide in enumerate(self.presentation.slides):
            slides.append(Slide(slide, index,slide_width= self.slide_width, slide_height= self.slide_height))
        return slides

    def find_shape_index(self,slide_index:int,shape_id:int)-> int:
        shapes = self.presentation.slides[slide_index].shapes

        shape_index = -1
        for _shape_index, shape in enumerate(shapes):
            if shape.shape_id == shape_id:
                return _shape_index
        return shape_index